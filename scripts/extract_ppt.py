import traceback

try:
    from pptx import Presentation

    prs = Presentation("c:/Users/ijara/Documentos/Codigos/python/coursera/Applied Data Science Capstone/ds-capstone-template-coursera.pptx")
    
    with open("c:/Users/ijara/Documentos/Codigos/python/coursera/Applied Data Science Capstone/ppt_structure.txt", "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text
            
            f.write(f"--- Slide {i+1} ---\n")
            f.write(f"Title: {title}\n")
            
            # Extract notes to know what Coursera wants on each slide (e.g. "Include a bar chart...")
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else "No notes"
            f.write(f"Notes: {notes}\n")
            
            # Get other text content
            content = []
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                if shape == slide.shapes.title: continue
                content.append(shape.text.replace("\n", " | "))
            f.write(f"Content: {content}\n\n")

    print("Success: ppt_structure.txt created.")
except ImportError:
    print("python-pptx not installed. Please install it.")
except Exception as e:
    print(f"Error: {e}")
    traceback.print_exc()
